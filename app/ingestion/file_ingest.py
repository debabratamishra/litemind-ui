"""Local-first document ingestion pipeline for RAG."""

from __future__ import annotations

import io
import json
import logging
import mimetypes
import re
import shutil
import subprocess
import tempfile
import tomllib
from pathlib import Path
from typing import Iterable, List, Optional, Tuple

from PIL import Image
from bs4 import BeautifulSoup

from app.core.rag_formats import ALLOWED_UPLOAD_EXTENSIONS, IMAGE_EXTENSIONS, LEGACY_OFFICE_EXTENSIONS, SUPPORTED_EXTENSIONS, TEXT_EXTENSIONS, TEXTISH_EXTENSIONS
from app.ingestion.enhanced_document_processor import extract_docx_enhanced, extract_epub_enhanced, extract_pdf_enhanced
from app.ingestion.enhanced_extractors import extract_csv_enhanced

logger = logging.getLogger(__name__)

try:
    import pandas as pd
    PANDAS_AVAILABLE = True
except ImportError:
    pd = None
    PANDAS_AVAILABLE = False

try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

try:
    import trafilatura
    TRAFILATURA_AVAILABLE = True
except ImportError:
    trafilatura = None
    TRAFILATURA_AVAILABLE = False

try:
    from striprtf.striprtf import rtf_to_text
    STRIPRTF_AVAILABLE = True
except ImportError:
    rtf_to_text = None
    STRIPRTF_AVAILABLE = False

try:
    from odf import teletype
    from odf.opendocument import load as load_odf_document
    ODFPY_AVAILABLE = True
except ImportError:
    teletype = None
    load_odf_document = None
    ODFPY_AVAILABLE = False

try:
    import yaml
    YAML_AVAILABLE = True
except ImportError:
    yaml = None
    YAML_AVAILABLE = False

try:
    from unstructured.partition.auto import partition
    from unstructured.documents.elements import Table
    UNSTRUCTURED_AVAILABLE = True
except ImportError:
    partition = None
    Table = None
    UNSTRUCTURED_AVAILABLE = False

DOT_ALLOWED_EXTENSIONS = ALLOWED_UPLOAD_EXTENSIONS
DOT_IMAGE_EXTENSIONS = {f".{ext}" for ext in IMAGE_EXTENSIONS}
DOT_TEXT_EXTENSIONS = {f".{ext}" for ext in TEXT_EXTENSIONS}
DOT_TEXTISH_EXTENSIONS = {f".{ext}" for ext in TEXTISH_EXTENSIONS}


def _guess_mime(path: Path) -> str:
    mime_type, _ = mimetypes.guess_type(str(path))
    return mime_type or "application/octet-stream"


def _read_file_bytes(path: Path) -> bytes:
    return path.read_bytes()


def _normalize_text(text: str) -> str:
    if not text:
        return ""
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    text = text.replace("\x00", "")
    text = re.sub(r"[ \t]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    text = re.sub(r"(?m)^ +", "", text)
    return text.strip()


def _make_block(content: str, path: Path, content_type: str, **metadata) -> Optional[dict]:
    clean_content = _normalize_text(content)
    if not clean_content:
        return None
    return {
        "content": clean_content,
        "metadata": {
            "filename": path.name,
            "filetype": _guess_mime(path),
            "content_type": content_type,
            **metadata,
        },
    }


def _split_logical_sections(text: str, max_chars: int = 2800) -> List[str]:
    clean_text = _normalize_text(text)
    if not clean_text:
        return []

    paragraphs = [part.strip() for part in re.split(r"\n\s*\n+", clean_text) if part.strip()]
    if not paragraphs:
        return [clean_text]

    sections: List[str] = []
    current: List[str] = []
    current_len = 0

    for paragraph in paragraphs:
        if len(paragraph) > max_chars:
            if current:
                sections.append("\n\n".join(current))
                current = []
                current_len = 0
            sections.extend(_split_long_text(paragraph, max_chars))
            continue

        projected_length = current_len + len(paragraph) + (2 if current else 0)
        if current and projected_length > max_chars:
            sections.append("\n\n".join(current))
            current = [paragraph]
            current_len = len(paragraph)
            continue

        current.append(paragraph)
        current_len = projected_length

    if current:
        sections.append("\n\n".join(current))

    return sections or [clean_text]


def _split_long_text(text: str, max_chars: int) -> List[str]:
    sentences = re.split(r"(?<=[.!?])\s+", text)
    chunks: List[str] = []
    current = ""

    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence:
            continue
        if len(sentence) > max_chars:
            if current:
                chunks.append(current.strip())
                current = ""
            for index in range(0, len(sentence), max_chars):
                piece = sentence[index:index + max_chars].strip()
                if piece:
                    chunks.append(piece)
            continue

        if current and len(current) + len(sentence) + 1 > max_chars:
            chunks.append(current.strip())
            current = sentence
            continue

        current = f"{current} {sentence}".strip() if current else sentence

    if current:
        chunks.append(current.strip())

    return chunks or [_normalize_text(text)]


def _collect_blocks_from_text(text: str, path: Path, content_type: str, **metadata) -> List[dict]:
    return [
        block
        for section in _split_logical_sections(text)
        if (block := _make_block(section, path, content_type, **metadata)) is not None
    ]


def _find_libreoffice_executable() -> Optional[str]:
    candidates = [
        shutil.which("soffice"),
        shutil.which("libreoffice"),
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
    ]
    for candidate in candidates:
        if candidate and Path(candidate).exists():
            return candidate
    return None


def _convert_with_libreoffice(path: Path, target_extension: str) -> Optional[Path]:
    soffice = _find_libreoffice_executable()
    if not soffice:
        return None

    out_dir = Path(tempfile.mkdtemp(prefix="litemind_convert_"))
    format_name = target_extension.lstrip(".")

    try:
        process = subprocess.run(
            [
                soffice,
                "--headless",
                "--convert-to",
                format_name,
                "--outdir",
                str(out_dir),
                str(path),
            ],
            capture_output=True,
            text=True,
            check=False,
            timeout=120,
        )
        if process.returncode != 0:
            logger.warning("LibreOffice conversion failed for %s: %s", path.name, process.stderr.strip())
            return None

        converted = out_dir / f"{path.stem}{target_extension}"
        if converted.exists():
            return converted
    except Exception as exc:
        logger.warning("LibreOffice conversion failed for %s: %s", path.name, exc)

    return None


def _extract_image_record(path: Path) -> Tuple[List[dict], List[dict], List[dict]]:
    try:
        image_bytes = _read_file_bytes(path)
        image = Image.open(io.BytesIO(image_bytes))
        image.verify()
    except Exception as exc:
        logger.warning("Invalid image file %s: %s", path.name, exc)
        return [], [], []

    return [], [
        {
            "image_bytes": image_bytes,
            "caption": f"Image file: {path.stem}",
            "metadata": {
                "filename": path.name,
                "filetype": _guess_mime(path),
            },
        }
    ], []


def _extract_html_document(path: Path) -> List[dict]:
    raw_html = path.read_text(encoding="utf-8", errors="ignore")
    extracted = None
    if TRAFILATURA_AVAILABLE:
        try:
            extracted = trafilatura.extract(raw_html, include_tables=True, include_comments=False)
        except Exception as exc:
            logger.debug("Trafilatura failed for %s: %s", path.name, exc)

    if extracted:
        return _collect_blocks_from_text(extracted, path, "html_content")

    soup = BeautifulSoup(raw_html, "html.parser")
    for tag_name in ("script", "style", "noscript"):
        for tag in soup.find_all(tag_name):
            tag.decompose()

    sections: List[dict] = []
    title = soup.title.get_text(" ", strip=True) if soup.title else ""
    if title:
        title_block = _make_block(title, path, "html_title", title=title)
        if title_block is not None:
            sections.append(title_block)

    headings = [heading.get_text(" ", strip=True) for heading in soup.find_all(["h1", "h2", "h3", "h4"])]
    body_text = soup.get_text("\n")
    blocks = _collect_blocks_from_text(body_text, path, "html_content", title=title, headings=headings[:10])
    return sections + blocks


def _extract_json_document(path: Path) -> List[dict]:
    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    try:
        parsed = json.loads(raw_text)
        pretty_text = json.dumps(parsed, indent=2, ensure_ascii=False, sort_keys=True)
        return _collect_blocks_from_text(pretty_text, path, "structured_text")
    except json.JSONDecodeError:
        return _collect_blocks_from_text(raw_text, path, "text_content")


def _extract_jsonl_document(path: Path) -> List[dict]:
    records = []
    for line in path.read_text(encoding="utf-8", errors="ignore").splitlines():
        line = line.strip()
        if not line:
            continue
        try:
            records.append(json.dumps(json.loads(line), ensure_ascii=False, sort_keys=True))
        except json.JSONDecodeError:
            records.append(line)
    return _collect_blocks_from_text("\n".join(records), path, "structured_text")


def _extract_yaml_document(path: Path) -> List[dict]:
    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    if not YAML_AVAILABLE:
        return _collect_blocks_from_text(raw_text, path, "text_content")

    try:
        parsed = yaml.safe_load(raw_text)
        pretty_text = json.dumps(parsed, indent=2, ensure_ascii=False, sort_keys=True)
        return _collect_blocks_from_text(pretty_text, path, "structured_text")
    except Exception:
        return _collect_blocks_from_text(raw_text, path, "text_content")


def _extract_toml_document(path: Path) -> List[dict]:
    raw_bytes = path.read_bytes()
    try:
        parsed = tomllib.loads(raw_bytes.decode("utf-8", errors="ignore"))
        pretty_text = json.dumps(parsed, indent=2, ensure_ascii=False, sort_keys=True)
        return _collect_blocks_from_text(pretty_text, path, "structured_text")
    except Exception:
        return _collect_blocks_from_text(raw_bytes.decode("utf-8", errors="ignore"), path, "text_content")


def _extract_rtf_document(path: Path) -> List[dict]:
    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    if STRIPRTF_AVAILABLE:
        try:
            raw_text = rtf_to_text(raw_text)
        except Exception as exc:
            logger.warning("RTF parsing failed for %s: %s", path.name, exc)
    return _collect_blocks_from_text(raw_text, path, "rtf_content")


def _extract_odt_document(path: Path) -> List[dict]:
    if not ODFPY_AVAILABLE:
        return _extract_with_unstructured_fallback(path)

    try:
        document = load_odf_document(str(path))
        text = teletype.extractText(document.text)
        return _collect_blocks_from_text(text, path, "odt_content")
    except Exception as exc:
        logger.warning("ODT parsing failed for %s: %s", path.name, exc)
        return _extract_with_unstructured_fallback(path)


def _extract_workbook_document(path: Path) -> List[dict]:
    if not PANDAS_AVAILABLE:
        logger.warning("Pandas is not available for spreadsheet extraction")
        return []

    blocks: List[dict] = []
    try:
        workbook = pd.ExcelFile(path)
        for sheet_name in workbook.sheet_names:
            try:
                dataframe = workbook.parse(sheet_name, dtype=str).fillna("")
            except Exception as exc:
                logger.warning("Failed to parse sheet %s in %s: %s", sheet_name, path.name, exc)
                continue

            if dataframe.empty:
                continue

            overview = _make_block(
                f"Workbook: {path.name}\nSheet: {sheet_name}\nRows: {len(dataframe)}\nColumns: {', '.join(str(column) for column in dataframe.columns[:20])}",
                path,
                "spreadsheet_overview",
                sheet=sheet_name,
                row_count=len(dataframe),
                column_count=len(dataframe.columns),
            )
            if overview is not None:
                blocks.append(overview)

            for start_index in range(0, min(len(dataframe), 1000), 75):
                chunk = dataframe.iloc[start_index:start_index + 75]
                header_line = " | ".join(str(column) for column in chunk.columns)
                row_lines = [header_line]
                for row_number, (_, row_values) in enumerate(chunk.iterrows(), start=start_index + 1):
                    row_lines.append(
                        f"Row {row_number}: " + " | ".join(str(value) for value in row_values.tolist())
                    )
                block = _make_block(
                    "\n".join(row_lines),
                    path,
                    "spreadsheet_rows",
                    sheet=sheet_name,
                    row_start=start_index + 1,
                    row_end=start_index + len(chunk),
                )
                if block is not None:
                    blocks.append(block)
    except Exception as exc:
        logger.warning("Spreadsheet extraction failed for %s: %s", path.name, exc)
    return blocks


def _extract_text_document(path: Path) -> List[dict]:
    suffix = path.suffix.lower()
    if suffix in {".html", ".htm", ".xml"}:
        return _extract_html_document(path)
    if suffix == ".json":
        return _extract_json_document(path)
    if suffix == ".jsonl":
        return _extract_jsonl_document(path)
    if suffix in {".yaml", ".yml"}:
        return _extract_yaml_document(path)
    if suffix == ".toml":
        return _extract_toml_document(path)

    raw_text = path.read_text(encoding="utf-8", errors="ignore")
    return _collect_blocks_from_text(raw_text, path, "text_content")


def _extract_presentation_document(path: Path) -> Tuple[List[dict], List[dict], List[dict]]:
    if not PPTX_AVAILABLE:
        return _extract_with_unstructured_fallback(path), [], []

    text_chunks: List[dict] = []
    image_records: List[dict] = []
    presentation = Presentation(str(path))

    for slide_index, slide in enumerate(presentation.slides, start=1):
        slide_parts: List[str] = []
        slide_title = ""

        title_shape = slide.shapes.title
        if title_shape is not None and getattr(title_shape, "text", "").strip():
            slide_title = title_shape.text.strip()
            slide_parts.append(f"Slide Title: {slide_title}")

        for shape_index, shape in enumerate(slide.shapes, start=1):
            try:
                if getattr(shape, "has_text_frame", False):
                    text = shape.text_frame.text.strip()
                    if text:
                        slide_parts.append(text)

                if getattr(shape, "has_table", False):
                    rows = []
                    for row in shape.table.rows:
                        cells = [cell.text_frame.text.strip() for cell in row.cells]
                        rows.append(" | ".join(cells))
                    if rows:
                        slide_parts.append("Table:\n" + "\n".join(rows))

                if hasattr(shape, "image"):
                    image_records.append(
                        {
                            "image_bytes": shape.image.blob,
                            "caption": f"{path.stem} - Slide {slide_index}, Image {shape_index}",
                            "metadata": {
                                "filename": path.name,
                                "slide_number": slide_index,
                                "shape_index": shape_index,
                                "filetype": _guess_mime(path),
                            },
                        }
                    )

                if getattr(shape, "has_chart", False):
                    chart = shape.chart
                    chart_title = ""
                    if getattr(chart, "has_title", False) and chart.chart_title is not None:
                        chart_title = chart.chart_title.text_frame.text.strip()
                    if chart_title:
                        slide_parts.append(f"Chart: {chart_title}")
            except Exception as exc:
                logger.debug("Failed to parse slide %s shape %s in %s: %s", slide_index, shape_index, path.name, exc)

        try:
            notes_frame = slide.notes_slide.notes_text_frame
            notes_text = notes_frame.text.strip()
            if notes_text:
                slide_parts.append(f"Speaker Notes:\n{notes_text}")
        except Exception:
            pass

        slide_text = "\n\n".join(part for part in slide_parts if part.strip())
        block = _make_block(
            slide_text,
            path,
            "presentation_slide",
            slide_number=slide_index,
            slide_title=slide_title,
        )
        if block is not None:
            text_chunks.append(block)

    return text_chunks, image_records, []


def _extract_with_unstructured_fallback(path: Path) -> List[dict]:
    if not UNSTRUCTURED_AVAILABLE:
        return []

    try:
        blocks: List[dict] = []
        tables: List[dict] = []
        for element in partition(filename=str(path)):
            text = getattr(element, "text", "") or ""
            if not text.strip():
                continue
            metadata = {
                "filename": path.name,
                "filetype": _guess_mime(path),
                "page_number": getattr(getattr(element, "metadata", None), "page_number", None),
            }
            if Table is not None and isinstance(element, Table):
                table_block = _make_block(text, path, "table_content", **metadata)
                if table_block is not None:
                    tables.append(table_block)
                continue

            block = _make_block(text, path, "unstructured_content", category=getattr(element, "category", None), **metadata)
            if block is not None:
                blocks.append(block)
        return blocks + tables
    except Exception as exc:
        logger.warning("Unstructured fallback failed for %s: %s", path.name, exc)
        return []


def _restore_original_metadata(records: List, original_path: Path) -> List:
    restored_records = []
    for record in records:
        if not isinstance(record, dict):
            restored_records.append(record)
            continue

        restored_records.append(
            {
                **record,
                "metadata": {
                    **record.get("metadata", {}),
                    "filename": original_path.name,
                    "source_filetype": _guess_mime(original_path),
                    "converted_from": original_path.suffix.lower(),
                },
            }
        )
    return restored_records


def _extract_legacy_office_document(path: Path) -> Tuple[List[dict], List[dict], List[dict]]:
    target_extension = {
        ".doc": ".docx",
        ".ppt": ".pptx",
        ".xls": ".xlsx",
    }.get(path.suffix.lower())
    if not target_extension:
        return [], [], []

    converted_path = _convert_with_libreoffice(path, target_extension)
    if converted_path is None:
        fallback_blocks = _extract_with_unstructured_fallback(path)
        return fallback_blocks, [], []

    try:
        text_chunks, image_records, table_texts = ingest_file(converted_path)
        return (
            _restore_original_metadata(text_chunks, path),
            _restore_original_metadata(image_records, path),
            _restore_original_metadata(table_texts, path),
        )
    finally:
        shutil.rmtree(converted_path.parent, ignore_errors=True)


def get_ingestion_capabilities() -> dict:
    return {
        "supported_extensions": SUPPORTED_EXTENSIONS,
        "local_pipeline": {
            "trafilatura": TRAFILATURA_AVAILABLE,
            "striprtf": STRIPRTF_AVAILABLE,
            "odfpy": ODFPY_AVAILABLE,
            "python_pptx": PPTX_AVAILABLE,
            "pandas": PANDAS_AVAILABLE,
            "yaml": YAML_AVAILABLE,
            "unstructured_fallback": UNSTRUCTURED_AVAILABLE,
            "libreoffice": _find_libreoffice_executable() is not None,
        },
    }


def ingest_file(path: Path) -> Tuple[List[dict], List[dict], List[dict]]:
    ext = path.suffix.lower()
    logger.info("Processing file %s with local-first pipeline", path.name)

    if ext not in DOT_ALLOWED_EXTENSIONS:
        logger.warning("Unsupported file format: %s", ext)
        return [], [], []

    if ext in DOT_IMAGE_EXTENSIONS:
        return _extract_image_record(path)

    if ext == ".pdf":
        return extract_pdf_enhanced(path), [], []

    if ext == ".docx":
        return extract_docx_enhanced(path), [], []

    if ext == ".epub":
        return extract_epub_enhanced(path), [], []

    if ext in {".csv", ".tsv"}:
        return extract_csv_enhanced(path), [], []

    if ext in {".xlsx"}:
        return _extract_workbook_document(path), [], []

    if ext in {".pptx"}:
        return _extract_presentation_document(path)

    if ext in {".doc", ".ppt", ".xls"}:
        return _extract_legacy_office_document(path)

    if ext == ".rtf":
        return _extract_rtf_document(path), [], []

    if ext == ".odt":
        return _extract_odt_document(path), [], []

    if ext in DOT_TEXTISH_EXTENSIONS or ext in {".html", ".htm", ".xml"}:
        return _extract_text_document(path), [], []

    fallback_chunks = _extract_with_unstructured_fallback(path)
    return fallback_chunks, [], []